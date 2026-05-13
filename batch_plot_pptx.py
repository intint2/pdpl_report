"""Batch matplotlib spectra per condition folder into one PowerPoint deck."""

from __future__ import annotations

import csv
from datetime import datetime
from io import BytesIO
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: E402

from pptx import Presentation
from pptx.chart.data import XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn, nsdecls
from pptx.util import Emu, Inches, Pt
from tqdm.auto import tqdm


def _short_postfix(name: str, max_len: int = 40) -> str:
    if len(name) <= max_len:
        return name
    return f"{name[: max_len - 1]}…"


# Matplotlib tab10 / Office-friendly distinct line colors
_TAB10_RGB: tuple[tuple[int, int, int], ...] = (
    (31, 119, 180),
    (255, 127, 14),
    (44, 160, 44),
    (214, 39, 40),
    (148, 103, 189),
    (140, 86, 75),
    (227, 119, 194),
    (127, 127, 127),
    (188, 189, 34),
    (23, 190, 207),
)

# Office charts: major gridlines (light gray ~ HTML lightgray)
_NATIVE_GRID_LINE_GRAY = RGBColor(211, 211, 211)
# Office OOXML: ~1 device px at 96 dpi (line width in EMU)
_NATIVE_1PX_LINE_EMU = 9525


def _apply_native_legend_frame(legend) -> None:
    """Legend: white fill, black border, ~1px outline (via c:spPr XML)."""
    if legend is None:
        return
    elm = legend._element
    old = elm.find(qn("c:spPr"))
    if old is not None:
        elm.remove(old)
    sp_pr = parse_xml(
        f'<c:spPr {nsdecls("a", "c")}>'
        '<a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>'
        f'<a:ln w="{_NATIVE_1PX_LINE_EMU}" cap="flat" cmpd="sng" algn="ctr">'
        '<a:solidFill><a:srgbClr val="000000"/></a:solidFill>'
        "</a:ln>"
        "</c:spPr>"
    )
    tx = elm.find(qn("c:txPr"))
    if tx is not None:
        elm.insert(list(elm).index(tx), sp_pr)
    else:
        elm.append(sp_pr)


def _apply_native_major_grid_gray(x_axis, y_axis) -> None:
    """Major gridlines: light gray (Office native charts only)."""
    for ax in (x_axis, y_axis):
        if not ax.has_major_gridlines:
            continue
        ln = ax.major_gridlines.format.line
        ln.color.rgb = _NATIVE_GRID_LINE_GRAY
        ln.width = Pt(0.75)


def _read_spectrum_csv(path: Path) -> tuple[list[float], list[float]]:
    w: list[float] = []
    i: list[float] = []
    with path.open(newline="", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        for row in reader:
            w.append(float(row["Wavelength_nm"]))
            i.append(float(row["Intensity"]))
    return w, i


def _legend_label(path: Path) -> str:
    stem = path.stem
    if "__" in stem:
        return stem.rsplit("__", maxsplit=1)[-1]
    return stem


def _apply_powd_xy_series_and_chrome(chart) -> None:
    """Distinct series colors, typography, legend placement, gallery style."""
    chart.chart_style = 12
    chart.font.size = Pt(10)
    legend = chart.legend
    if legend is not None:
        legend.position = XL_LEGEND_POSITION.BOTTOM
        legend.include_in_layout = False
        legend.font.size = Pt(8)
    for idx, ser in enumerate(chart.series):
        r, g, b = _TAB10_RGB[idx % len(_TAB10_RGB)]
        ser.format.line.color.rgb = RGBColor(r, g, b)
        ser.format.line.width = Emu(_NATIVE_1PX_LINE_EMU)


def _condition_figure_png(cond_dir: Path) -> bytes | None:
    files = sorted(cond_dir.glob("*.csv"), key=lambda p: p.name.casefold())
    if not files:
        return None
    fig, ax = plt.subplots(figsize=(10, 5.625))
    for idx, path in enumerate(files):
        wl, intensity = _read_spectrum_csv(path)
        r, g, b = (c / 255.0 for c in _TAB10_RGB[idx % len(_TAB10_RGB)])
        ax.plot(
            wl,
            intensity,
            label=_legend_label(path),
            linewidth=1.05,
            color=(r, g, b),
        )
    ax.set_xlabel("Wavelength (nm)")
    ax.set_ylabel("Intensity")
    ax.set_title(cond_dir.name)
    ax.set_ylim(bottom=0)
    ax.set_axisbelow(True)
    ax.grid(True, color="lightgray", linestyle="--", linewidth=0.8)
    ax.tick_params(axis="both", which="both", direction="in", top=True, right=True)
    ax.legend(fontsize=8, loc="upper right", ncol=2, framealpha=0.92)
    fig.tight_layout()
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=120)
    plt.close(fig)
    return buf.getvalue()


def _default_output_dir(project_root: Path) -> Path:
    return project_root / "output"


def run_powd_batch_report(
    powd_root: Path,
    *,
    project_root: Path,
    output_dir: Path | None = None,
) -> Path:
    """
    One title slide plus one image slide per condition subfolder (all CSV spectra overlay).
    Saves ``powd_pl_report_YYYYMMDD_HHMMSS.pptx`` under output_dir.
    """
    root = powd_root.resolve()
    out_dir = output_dir if output_dir is not None else _default_output_dir(project_root)
    out_dir.mkdir(parents=True, exist_ok=True)
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    out_path = out_dir / f"powd_pl_report_{stamp}.pptx"

    folders = sorted((p for p in root.iterdir() if p.is_dir()), key=lambda p: p.name.casefold())
    total_folders = len(folders)

    prs = Presentation()
    layout_title = prs.slide_layouts[0]
    layout_blank = prs.slide_layouts[6]
    layout_content = prs.slide_layouts[1]

    s0 = prs.slides.add_slide(layout_title)
    s0.shapes.title.text = "POWD PL Report"
    try:
        s0.placeholders[1].text = str(root)
    except (IndexError, KeyError, AttributeError):
        pass

    plots_added = 0
    bar = tqdm(
        folders,
        desc="Plot / slides",
        unit="folder",
        total=total_folders,
        dynamic_ncols=True,
    )
    for cond in bar:
        bar.set_postfix_str(_short_postfix(cond.name), refresh=False)
        png = _condition_figure_png(cond)
        if png is None:
            tqdm.write(f"Skipped (no CSV): {cond.name}")
            continue
        slide = prs.slides.add_slide(layout_blank)
        stream = BytesIO(png)
        slide.shapes.add_picture(stream, Inches(0.4), Inches(0.35), width=Inches(9.2))
        plots_added += 1
        bar.set_postfix_str(f"{_short_postfix(cond.name, 28)} | slides {plots_added}", refresh=True)

    if plots_added == 0:
        s1 = prs.slides.add_slide(layout_content)
        s1.shapes.title.text = "No plotable CSV data"
        s1.placeholders[1].text = f"No *.csv files found under condition folders in:\n{root}"

    tqdm.write("Saving PPTX…")
    prs.save(out_path)
    tqdm.write(f"Done: {plots_added} image slide(s) → {out_path}")
    return out_path


def run_powd_batch_report_editable(
    powd_root: Path,
    *,
    project_root: Path,
    output_dir: Path | None = None,
) -> Path:
    """
    Same layout as :func:`run_powd_batch_report`, but each condition uses a native
    PowerPoint **line (XY) chart** (embedded Excel data). Open in Office and use
    **Chart Design** or right-click **Edit Data** to adjust series, colors, axes, etc.
    """
    root = powd_root.resolve()
    out_dir = output_dir if output_dir is not None else _default_output_dir(project_root)
    out_dir.mkdir(parents=True, exist_ok=True)
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    out_path = out_dir / f"powd_pl_report_edit_{stamp}.pptx"

    folders = sorted((p for p in root.iterdir() if p.is_dir()), key=lambda p: p.name.casefold())
    total_folders = len(folders)

    prs = Presentation()
    layout_title = prs.slide_layouts[0]
    layout_blank = prs.slide_layouts[6]
    layout_content = prs.slide_layouts[1]

    s0 = prs.slides.add_slide(layout_title)
    s0.shapes.title.text = "POWD PL Report (editable charts)"
    try:
        s0.placeholders[1].text = str(root)
    except (IndexError, KeyError, AttributeError):
        pass

    plots_added = 0
    bar = tqdm(
        folders,
        desc="Charts / slides",
        unit="folder",
        total=total_folders,
        dynamic_ncols=True,
    )
    for cond in bar:
        bar.set_postfix_str(_short_postfix(cond.name), refresh=False)
        files = sorted(cond.glob("*.csv"), key=lambda p: p.name.casefold())
        if not files:
            tqdm.write(f"Skipped (no CSV): {cond.name}")
            continue

        chart_data = XyChartData()
        x_min = float("inf")
        x_max = float("-inf")
        for path in files:
            wl, intensity = _read_spectrum_csv(path)
            x_min = min(x_min, min(wl))
            x_max = max(x_max, max(wl))
            series = chart_data.add_series(_legend_label(path))
            for x, y in zip(wl, intensity, strict=True):
                series.add_data_point(x, y)

        slide = prs.slides.add_slide(layout_blank)
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS,
            Inches(0.45),
            Inches(0.55),
            Inches(9.1),
            Inches(5.45),
            chart_data,
        )
        chart = chart_shape.chart
        chart.has_legend = True
        chart.has_title = True
        chart.chart_title.text_frame.text = cond.name

        xa = chart.category_axis
        xa.has_title = True
        xa.axis_title.text_frame.text = "Wavelength (nm)"
        xa.has_major_gridlines = True
        xa.major_tick_mark = XL_TICK_MARK.INSIDE
        xa.minor_tick_mark = XL_TICK_MARK.INSIDE

        ya = chart.value_axis
        ya.has_title = True
        ya.axis_title.text_frame.text = "Intensity"
        ya.has_major_gridlines = True
        ya.major_tick_mark = XL_TICK_MARK.INSIDE
        ya.minor_tick_mark = XL_TICK_MARK.INSIDE
        ya.minimum_scale = 0.0

        xa.minimum_scale = float(round(x_min))
        xa.maximum_scale = float(round(x_max))
        xa.tick_labels.number_format = "0"

        _apply_powd_xy_series_and_chrome(chart)
        _apply_native_major_grid_gray(xa, ya)
        _apply_native_legend_frame(chart.legend)

        plots_added += 1
        bar.set_postfix_str(f"{_short_postfix(cond.name, 28)} | slides {plots_added}", refresh=True)

    if plots_added == 0:
        s1 = prs.slides.add_slide(layout_content)
        s1.shapes.title.text = "No plotable CSV data"
        s1.placeholders[1].text = f"No *.csv files found under condition folders in:\n{root}"

    tqdm.write("Saving PPTX…")
    prs.save(out_path)
    tqdm.write(f"Done: {plots_added} chart slide(s) → {out_path}")
    return out_path
